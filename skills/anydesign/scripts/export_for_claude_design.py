#!/usr/bin/env python3
"""
export_for_claude_design.py — Bridge anydesign output to Claude Design.

Takes a `design.md` (+ optional `design-tokens.json`) and emits a bundle folder
ready to upload as "brand and product assets" during Claude Design setup.

Outputs (each can be skipped with --skip):
  tokens.css              CSS custom properties from DTCG tokens
  tailwind.config.ts      Tailwind v3 config from DTCG tokens
  brand-overview.docx     The design.md body, rendered as Word doc
  brand-kit.pptx          Cover + palette + typography + components + Do's/Don'ts deck
  README-claude-design.md Upload instructions for the user

Usage:
    python export_for_claude_design.py design.md
    python export_for_claude_design.py design.md design-tokens.json --out my-bundle/
    python export_for_claude_design.py design.md --skip pptx,docx     # CSS/Tailwind only

Optional deps (script gracefully skips formats whose dep is missing):
    pyyaml          — parse YAML frontmatter (needed for everything except tokens.css/tailwind)
    python-pptx     — brand-kit.pptx
    python-docx     — brand-overview.docx
"""

import argparse
import json
import re
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
if hasattr(sys.stderr, "reconfigure"):
    try:
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Inches as DocxInches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


FRONTMATTER_RE = re.compile(r"\A---\s*\n(.*?)\n---\s*\n", re.DOTALL)
HEX_RE = re.compile(r"^#[0-9a-fA-F]{3,8}$")
ALL_FORMATS = ("css", "tailwind", "docx", "pptx", "readme")


# =============================================================================
# design.md parsing
# =============================================================================

def split_frontmatter(text):
    m = FRONTMATTER_RE.match(text)
    if not m:
        return None, text
    return m.group(1), text[m.end():]


def parse_frontmatter(fm_text):
    if not HAS_YAML:
        return None
    try:
        return yaml.safe_load(fm_text)
    except yaml.YAMLError as e:
        print(f"  warn: frontmatter YAML parse error: {e}", file=sys.stderr)
        return None


def extract_section(body, section_number):
    pattern = rf"^##\s+{section_number}\.\s.*?\n(.*?)(?=^##\s+\d+\.\s|\Z)"
    m = re.search(pattern, body, re.DOTALL | re.MULTILINE)
    return m.group(1).strip() if m else None


def extract_dos_donts(body):
    section6 = extract_section(body, 6)
    if not section6:
        return [], []
    do_match = re.search(r"###\s+Do\s*\n(.*?)(?=###\s+Don|\Z)", section6, re.DOTALL)
    dont_match = re.search(r"###\s+Don[''`]t\s*\n(.*?)\Z", section6, re.DOTALL)

    def parse_bullets(text):
        if not text:
            return []
        bullets = []
        current = None
        for line in text.split("\n"):
            stripped = line.strip()
            if re.match(r"^[-*]\s+", stripped):
                if current:
                    bullets.append(current.strip())
                current = re.sub(r"^[-*]\s+", "", stripped)
            elif current and stripped and not stripped.startswith("#"):
                current += " " + stripped
        if current:
            bullets.append(current.strip())
        return bullets

    return (
        parse_bullets(do_match.group(1) if do_match else ""),
        parse_bullets(dont_match.group(1) if dont_match else ""),
    )


def extract_tldr(body):
    m = re.search(r"^##\s+TL;DR\s*\n(.*?)(?=^##\s+|\Z)", body, re.DOTALL | re.MULTILINE)
    if not m:
        return ""
    text = m.group(1).strip()
    text = re.sub(r"^---\s*$", "", text, flags=re.MULTILINE).strip()
    return text


# =============================================================================
# DTCG token tree
# =============================================================================

def is_token_leaf(node):
    return isinstance(node, dict) and "$value" in node


def walk_tokens(tree, prefix=()):
    """Yield (path_tuple, value, type) for every DTCG leaf."""
    if not isinstance(tree, dict):
        return
    for key, val in tree.items():
        if key.startswith("$"):
            continue
        path = prefix + (key,)
        if is_token_leaf(val):
            yield path, val["$value"], val.get("$type", "unknown")
        elif isinstance(val, dict):
            yield from walk_tokens(val, path)


def collect_by_type(tree, *types):
    out = []
    for path, value, t in walk_tokens(tree):
        if t in types:
            out.append((path, value, t))
    return out


def slug(path):
    return "-".join(str(p) for p in path)


# =============================================================================
# CJK font fallback (issue #3)
# =============================================================================

# Noto * SC first: available via Google Fonts, so cloud render environments
# (claude.ai/design) can resolve it; PingFang/Hiragino/YaHei cover local macOS/Windows.
CJK_FALLBACKS = {
    "sans-serif": ["Noto Sans SC", "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei"],
    "serif": ["Noto Serif SC", "Songti SC", "SimSun"],
    "monospace": ["Noto Sans Mono CJK SC", "PingFang SC"],
}

CJK_FONT_RE = re.compile(
    r"noto\s+(sans|serif)\s+(sc|tc|jp|kr)|noto\s+sans\s+mono\s+cjk|pingfang|hiragino\s+sans\s+gb"
    r"|microsoft\s+(yahei|jhenghei)|source\s+han|simsun|simhei|songti|heiti|wenquanyi|lxgw",
    re.IGNORECASE,
)

MONO_HINT_RE = re.compile(r"mono|consolas|menlo|monaco|courier|sfmono", re.IGNORECASE)


def split_font_stack(value):
    """Normalize a DTCG fontFamily $value (string or list) to a list of names."""
    if isinstance(value, str):
        return [f.strip() for f in value.split(",") if f.strip()]
    if isinstance(value, list):
        return [str(f).strip() for f in value if str(f).strip()]
    return []


def _bare(name):
    return name.strip().strip("'\"").lower()


def classify_font_stack(families):
    bare = [_bare(f) for f in families]
    for f in bare:
        if f in ("monospace", "ui-monospace") or MONO_HINT_RE.search(f):
            return "monospace"
    for f in bare:
        if f in ("serif", "ui-serif"):
            return "serif"
    return "sans-serif"


def ensure_cjk_fallback(families):
    """Append a Simplified-Chinese fallback chain to a font stack, keeping the
    generic family (sans-serif/serif/monospace) last. No-op if the stack
    already names a CJK font."""
    if not families:
        return families
    if any(CJK_FONT_RE.search(f) for f in families):
        return list(families)
    generic = classify_font_stack(families)
    out = list(families)
    tail = out.pop() if _bare(out[-1]) == generic else None
    out += CJK_FALLBACKS[generic]
    out.append(tail if tail is not None else generic)
    return out


def css_font_name(name):
    name = name.strip()
    if name.startswith(("'", '"')) or " " not in name:
        return name
    return f'"{name}"'


def render_font_family_css(value, cjk_fallback=True):
    families = split_font_stack(value)
    if cjk_fallback:
        families = ensure_cjk_fallback(families)
    return ", ".join(css_font_name(f) for f in families)


# =============================================================================
# tokens.css
# =============================================================================

def emit_tokens_css(tokens, out_path, source_name, cjk_fallback=True):
    if not tokens:
        return False
    lines = [
        "/*",
        f" * Design tokens for {source_name}",
        " * Generated by anydesign · export_for_claude_design.py",
        " * Source: DTCG design-tokens.json",
        " */",
        "",
        ":root {",
    ]
    for path, value, ttype in walk_tokens(tokens):
        name = "--" + slug(path)
        if ttype == "fontFamily":
            rendered = render_font_family_css(value, cjk_fallback)
        else:
            rendered = render_css_value(value, ttype)
        lines.append(f"  {name}: {rendered};")
    lines.append("}")
    lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    return True


def render_css_value(value, ttype):
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, str):
        return value
    return json.dumps(value)


# =============================================================================
# tailwind.config.ts
# =============================================================================

def emit_tailwind_config(tokens, out_path, source_name, cjk_fallback=True):
    if not tokens:
        return False

    color = build_tailwind_colors(tokens.get("color", {}))
    typography = tokens.get("typography", {})
    font_family = build_tailwind_font_family(typography.get("font-family", {}), cjk_fallback)
    font_size = build_tailwind_scalar_map(typography.get("font-size", {}))
    font_weight = build_tailwind_scalar_map(typography.get("font-weight", {}))
    line_height = build_tailwind_scalar_map(typography.get("line-height", {}))
    letter_spacing = build_tailwind_scalar_map(typography.get("letter-spacing", {}))
    spacing = build_tailwind_scalar_map(tokens.get("spacing", {}))
    border_radius = build_tailwind_scalar_map(tokens.get("radius", {}))
    box_shadow = build_tailwind_scalar_map(tokens.get("shadow", {}))

    def js_obj(d, indent=4):
        if not d:
            return "{}"
        pad = " " * indent
        inner = ",\n".join(f"{pad}{json.dumps(k)}: {render_js(v, indent + 2)}" for k, v in d.items())
        return "{\n" + inner + "\n" + " " * (indent - 2) + "}"

    content = (
        "// tailwind.config.ts\n"
        f"// Generated by anydesign for {source_name}\n"
        "// Source: DTCG design-tokens.json\n\n"
        "import type { Config } from 'tailwindcss';\n\n"
        "const config: Config = {\n"
        "  content: ['./src/**/*.{html,js,jsx,ts,tsx}'],\n"
        "  theme: {\n"
        "    extend: {\n"
        f"      colors: {js_obj(color, 8)},\n"
        f"      fontFamily: {js_obj(font_family, 8)},\n"
        f"      fontSize: {js_obj(font_size, 8)},\n"
        f"      fontWeight: {js_obj(font_weight, 8)},\n"
        f"      lineHeight: {js_obj(line_height, 8)},\n"
        f"      letterSpacing: {js_obj(letter_spacing, 8)},\n"
        f"      spacing: {js_obj(spacing, 8)},\n"
        f"      borderRadius: {js_obj(border_radius, 8)},\n"
        f"      boxShadow: {js_obj(box_shadow, 8)},\n"
        "    },\n"
        "  },\n"
        "};\n\n"
        "export default config;\n"
    )
    out_path.write_text(content, encoding="utf-8")
    return True


def render_js(v, indent):
    if isinstance(v, dict):
        if not v:
            return "{}"
        pad = " " * indent
        inner = ",\n".join(f"{pad}{json.dumps(k)}: {render_js(val, indent + 2)}" for k, val in v.items())
        return "{\n" + inner + "\n" + " " * (indent - 2) + "}"
    if isinstance(v, (int, float)):
        return str(v)
    return json.dumps(v)


def build_tailwind_colors(color_tree):
    """Convert the DTCG color subtree into Tailwind's nested colors object."""
    out = {}
    for key, val in color_tree.items():
        if key.startswith("$"):
            continue
        if is_token_leaf(val):
            out[key] = val["$value"]
        elif isinstance(val, dict):
            nested = {}
            for sub_key, sub_val in val.items():
                if sub_key.startswith("$"):
                    continue
                if is_token_leaf(sub_val):
                    nested[sub_key] = sub_val["$value"]
                elif isinstance(sub_val, dict):
                    nested[sub_key] = build_tailwind_colors({sub_key: sub_val}).get(sub_key, {})
            if nested:
                out[key] = nested
    return out


def build_tailwind_scalar_map(subtree):
    """Flatten a one-level DTCG subtree to {name: value}."""
    out = {}
    for key, val in subtree.items():
        if key.startswith("$"):
            continue
        if is_token_leaf(val):
            v = val["$value"]
            if isinstance(v, (int, float)):
                out[key] = str(v)
            else:
                out[key] = v
    return out


def build_tailwind_font_family(subtree, cjk_fallback=True):
    out = {}
    for key, val in subtree.items():
        if key.startswith("$"):
            continue
        if is_token_leaf(val):
            families = split_font_stack(val["$value"])
            if cjk_fallback:
                families = ensure_cjk_fallback(families)
            out[key] = families
    return out


# =============================================================================
# brand-overview.docx
# =============================================================================

def emit_brand_overview_docx(fm, body, out_path, source_name):
    if not HAS_DOCX:
        return False

    doc = docx.Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    title = doc.add_heading(f"Brand overview — {source_name}", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    if fm and fm.get("description"):
        p = doc.add_paragraph()
        run = p.add_run(fm["description"].strip())
        run.italic = True

    if fm:
        meta = []
        if fm.get("source"):
            meta.append(f"Source: {fm['source']}")
        if fm.get("captured_at"):
            meta.append(f"Captured: {fm['captured_at']}")
        if meta:
            mp = doc.add_paragraph()
            mr = mp.add_run(" · ".join(meta))
            mr.font.size = DocxPt(9)
            mr.font.color.rgb = DocxRGB(0x6B, 0x72, 0x80)

    render_markdown_to_docx(body, doc)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    return True


def render_markdown_to_docx(body, doc):
    """Minimal markdown → docx renderer. Handles headings, bullets, tables, paragraphs."""
    lines = body.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
            i += 1
            continue
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=1)
            i += 1
            continue
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=2)
            i += 1
            continue
        if stripped.startswith("#### "):
            doc.add_heading(stripped[5:], level=3)
            i += 1
            continue

        if stripped == "---":
            i += 1
            continue

        if stripped.startswith(">"):
            p = doc.add_paragraph(style="Intense Quote") if "Intense Quote" in [s.name for s in doc.styles] else doc.add_paragraph()
            p.add_run(stripped.lstrip("> ").strip())
            i += 1
            continue

        if "|" in stripped and i + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-+:?\s*\|", lines[i + 1]):
            table_lines = []
            while i < len(lines) and "|" in lines[i]:
                table_lines.append(lines[i])
                i += 1
            render_markdown_table(table_lines, doc)
            continue

        if re.match(r"^[-*]\s+", stripped):
            doc.add_paragraph(re.sub(r"^[-*]\s+", "", stripped), style="List Bullet")
            i += 1
            continue

        if re.match(r"^\d+\.\s+", stripped):
            doc.add_paragraph(re.sub(r"^\d+\.\s+", "", stripped), style="List Number")
            i += 1
            continue

        para_lines = [stripped]
        i += 1
        while i < len(lines) and lines[i].strip() and not is_block_start(lines[i].strip()):
            para_lines.append(lines[i].strip())
            i += 1
        add_paragraph_with_inline(doc, " ".join(para_lines))


def is_block_start(line):
    if line.startswith(("#", ">", "- ", "* ", "---")):
        return True
    if re.match(r"^\d+\.\s", line):
        return True
    return False


def render_markdown_table(lines, doc):
    rows = []
    for ln in lines:
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        if cells and re.match(r"^:?-+:?$", cells[0]):
            continue
        rows.append(cells)
    if not rows:
        return
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]
    table = doc.add_table(rows=len(rows), cols=ncols)
    table.style = "Light Grid Accent 1"
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.rows[ri].cells[ci]
            cell.text = strip_inline_md(cell_text)
            if ri == 0:
                for p in cell.paragraphs:
                    for r in p.runs:
                        r.bold = True


def strip_inline_md(s):
    s = re.sub(r"`([^`]+)`", r"\1", s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", s)
    return s


def add_paragraph_with_inline(doc, text):
    p = doc.add_paragraph()
    parts = re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            run = p.add_run(part[2:-2])
            run.bold = True
        elif part.startswith("`") and part.endswith("`"):
            run = p.add_run(part[1:-1])
            run.font.name = "Consolas"
        else:
            p.add_run(part)


# =============================================================================
# brand-kit.pptx
# =============================================================================

SLIDE_W = Inches(13.333) if HAS_PPTX else None
SLIDE_H = Inches(7.5) if HAS_PPTX else None


def emit_brand_kit_pptx(fm, tokens, body, out_path, source_name):
    if not HAS_PPTX:
        return False

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover_slide(prs, fm, source_name)
    add_atmosphere_slide(prs, fm, body)
    add_palette_slides(prs, fm, tokens)
    add_typography_slide(prs, fm, tokens)
    add_spacing_radius_slide(prs, fm, tokens)
    add_components_slide(prs, fm)
    add_dos_donts_slide(prs, body)
    add_reconstruction_slide(prs, body)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return True


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text_box(slide, left, top, width, height, text, *, size=18, bold=False, color=None, align=None, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if font:
        r.font.name = font
    if color:
        r.font.color.rgb = RGBColor(*color)
    return tb


def hex_to_rgb(hexstr):
    h = hexstr.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) == 8:
        h = h[:6]
    if len(h) != 6:
        return None
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def add_cover_slide(prs, fm, source_name):
    s = add_blank_slide(prs)
    name = (fm or {}).get("name") or source_name
    source = (fm or {}).get("source", "")
    captured = (fm or {}).get("captured_at", "")

    add_text_box(s, Inches(0.6), Inches(2.4), Inches(12.0), Inches(1.5),
                 "Brand kit", size=20, color=(128, 128, 128))
    add_text_box(s, Inches(0.6), Inches(2.8), Inches(12.0), Inches(2.0),
                 name, size=54, bold=True, color=(23, 23, 23))
    if source:
        add_text_box(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(0.5),
                     source, size=14, color=(75, 85, 99), font="Consolas")
    if captured:
        add_text_box(s, Inches(0.6), Inches(5.5), Inches(12.0), Inches(0.5),
                     f"Captured {captured}", size=11, color=(156, 163, 175))
    add_text_box(s, Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.4),
                 "Generated with anydesign · for Claude Design",
                 size=9, color=(156, 163, 175))


def add_atmosphere_slide(prs, fm, body):
    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
                 "Atmosphere", size=28, bold=True, color=(23, 23, 23))

    text = ""
    if fm and fm.get("description"):
        text = fm["description"].strip()
    else:
        text = extract_tldr(body)

    if text:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for para in [p.strip() for p in text.split("\n\n") if p.strip()]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = para
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(12)


def add_palette_slides(prs, fm, tokens):
    swatches = collect_swatches(fm, tokens)
    if not swatches:
        return
    per_slide = 18
    chunks = [swatches[i:i + per_slide] for i in range(0, len(swatches), per_slide)]
    total = len(chunks)
    for idx, chunk in enumerate(chunks, 1):
        s = add_blank_slide(prs)
        title = "Color palette" if total == 1 else f"Color palette ({idx}/{total})"
        add_text_box(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
                     title, size=28, bold=True, color=(23, 23, 23))
        draw_swatch_grid(s, chunk)


def collect_swatches(fm, tokens):
    swatches = []
    seen = set()

    if tokens:
        for path, value, ttype in walk_tokens(tokens):
            if ttype != "color":
                continue
            if not isinstance(value, str):
                continue
            rgb = hex_to_rgb(value)
            if not rgb:
                continue
            key = (slug(path), value.upper())
            if key in seen:
                continue
            seen.add(key)
            swatches.append({"name": slug(path), "hex": value.upper(), "rgb": rgb})

    if not swatches and fm:
        colors = fm.get("colors") or {}
        for name, value in colors.items():
            if not isinstance(value, str):
                continue
            rgb = hex_to_rgb(value)
            if not rgb:
                continue
            swatches.append({"name": name, "hex": value.upper(), "rgb": rgb})

    return swatches


def draw_swatch_grid(slide, swatches):
    cols = 6
    margin_left = Inches(0.6)
    margin_top = Inches(1.5)
    cell_w = Inches(2.0)
    cell_h = Inches(1.7)
    gap_x = Inches(0.05)
    gap_y = Inches(0.1)

    for i, sw in enumerate(swatches):
        col = i % cols
        row = i // cols
        left = margin_left + col * (cell_w + gap_x)
        top = margin_top + row * (cell_h + gap_y)

        swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cell_w, Inches(1.1))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = RGBColor(*sw["rgb"])
        swatch.line.color.rgb = RGBColor(235, 235, 235)
        swatch.line.width = Pt(0.5)
        swatch.shadow.inherit = False

        label = slide.shapes.add_textbox(left, top + Inches(1.15), cell_w, Inches(0.5))
        tf = label.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = sw["name"]
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = RGBColor(23, 23, 23)

        hex_p = tf.add_paragraph()
        hex_run = hex_p.add_run()
        hex_run.text = sw["hex"]
        hex_run.font.size = Pt(8)
        hex_run.font.name = "Consolas"
        hex_run.font.color.rgb = RGBColor(107, 114, 128)


def add_typography_slide(prs, fm, tokens):
    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
                 "Typography", size=28, bold=True, color=(23, 23, 23))

    typography = (fm or {}).get("typography") or {}
    if not typography:
        return

    top = Inches(1.6)
    for token_name, style in typography.items():
        if not isinstance(style, dict):
            continue
        family = style.get("fontFamily", "")
        size = style.get("fontSize", "")
        weight = style.get("fontWeight", "")

        size_pt = parse_pt(size) or 16
        weight_int = int(weight) if weight and str(weight).isdigit() else 400
        bold = weight_int >= 600

        sample_box = s.shapes.add_textbox(Inches(0.6), top, Inches(8.0), Inches(0.7))
        sf = sample_box.text_frame
        sf.margin_top = Inches(0.0)
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = "The quick brown fox"
        sr.font.size = Pt(min(size_pt, 32))
        sr.font.bold = bold
        if family:
            sr.font.name = family.split(",")[0].strip().strip("'").strip('"')
        sr.font.color.rgb = RGBColor(23, 23, 23)

        meta_box = s.shapes.add_textbox(Inches(8.8), top + Inches(0.1), Inches(4.0), Inches(0.6))
        mf = meta_box.text_frame
        mp = mf.paragraphs[0]
        mr = mp.add_run()
        mr.text = token_name
        mr.font.size = Pt(11)
        mr.font.bold = True
        mr.font.color.rgb = RGBColor(23, 23, 23)
        mp2 = mf.add_paragraph()
        mr2 = mp2.add_run()
        mr2.text = f"{size or '?'} · {weight or '?'}"
        mr2.font.size = Pt(9)
        mr2.font.name = "Consolas"
        mr2.font.color.rgb = RGBColor(107, 114, 128)

        top += Inches(0.8)
        if top > Inches(7.0):
            break


def parse_pt(size_str):
    if not size_str:
        return None
    m = re.match(r"^([\d.]+)\s*(px|pt|rem|em)?$", str(size_str).strip())
    if not m:
        return None
    val = float(m.group(1))
    unit = (m.group(2) or "px").lower()
    if unit == "px":
        return val * 0.75
    if unit == "pt":
        return val
    if unit in ("rem", "em"):
        return val * 12
    return val


def add_spacing_radius_slide(prs, fm, tokens):
    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(6.0), Inches(0.6),
                 "Spacing scale", size=24, bold=True, color=(23, 23, 23))
    add_text_box(s, Inches(7.0), Inches(0.5), Inches(6.0), Inches(0.6),
                 "Radii", size=24, bold=True, color=(23, 23, 23))

    spacing_scale = []
    if fm and isinstance(fm.get("spacing"), dict):
        scale = fm["spacing"].get("scale")
        if isinstance(scale, list):
            spacing_scale = [(str(v), int(v) if isinstance(v, int) else parse_px(str(v))) for v in scale]
    if not spacing_scale and tokens:
        for path, value, ttype in walk_tokens(tokens.get("spacing", {})):
            px = parse_px(str(value))
            if px:
                spacing_scale.append((slug(path) + f" · {value}", px))

    top = Inches(1.5)
    for label, px in spacing_scale[:12]:
        bar_w = max(Inches(0.05), Inches(min(px / 32.0, 5.0)))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, bar_w, Inches(0.25))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(23, 23, 23)
        bar.line.fill.background()
        bar.shadow.inherit = False

        lbl = s.shapes.add_textbox(Inches(0.6) + bar_w + Inches(0.1), top - Inches(0.02),
                                   Inches(3.0), Inches(0.3))
        p = lbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(label)
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(75, 85, 99)
        top += Inches(0.4)

    radii = {}
    if fm and isinstance(fm.get("rounded"), dict):
        radii = {k: v for k, v in fm["rounded"].items() if isinstance(v, (str, int))}
    elif tokens:
        for path, value, ttype in walk_tokens(tokens.get("radius", {})):
            radii[slug(path)] = value

    top = Inches(1.5)
    for name, value in list(radii.items())[:10]:
        px = parse_px(str(value)) or 6
        display_px = min(px, 60)
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), top, Inches(1.2), Inches(0.6))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(245, 245, 245)
        chip.line.color.rgb = RGBColor(235, 235, 235)
        chip.adjustments[0] = min(display_px / 60.0 * 0.5, 0.5)
        chip.shadow.inherit = False

        lbl = s.shapes.add_textbox(Inches(8.3), top + Inches(0.05), Inches(4.0), Inches(0.5))
        p = lbl.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"{name} · {value}"
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(75, 85, 99)
        top += Inches(0.75)


def parse_px(value):
    m = re.match(r"^\s*([\d.]+)\s*(px|rem|em)?\s*$", str(value))
    if not m:
        return None
    val = float(m.group(1))
    unit = (m.group(2) or "px").lower()
    if unit == "px":
        return val
    if unit in ("rem", "em"):
        return val * 16
    return val


def add_components_slide(prs, fm):
    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
                 "Components", size=28, bold=True, color=(23, 23, 23))

    components = (fm or {}).get("components") or {}
    if not components:
        return

    top = Inches(1.5)
    col_left = [Inches(0.6), Inches(6.9)]
    col_top = [top, top]
    col_idx = 0

    for name, props in components.items():
        current_top = col_top[col_idx]
        if current_top > Inches(7.0):
            col_idx = 1
            current_top = col_top[col_idx]
        if current_top > Inches(7.0):
            break
        left = col_left[col_idx]

        name_box = s.shapes.add_textbox(left, current_top, Inches(6.0), Inches(0.4))
        np = name_box.text_frame.paragraphs[0]
        nr = np.add_run()
        nr.text = name
        nr.font.size = Pt(13)
        nr.font.bold = True
        nr.font.color.rgb = RGBColor(23, 23, 23)
        nr.font.name = "Consolas"

        desc = describe_component(props)
        desc_box = s.shapes.add_textbox(left, current_top + Inches(0.35),
                                        Inches(6.0), Inches(0.6))
        dp = desc_box.text_frame.paragraphs[0]
        desc_box.text_frame.word_wrap = True
        dr = dp.add_run()
        dr.text = desc
        dr.font.size = Pt(10)
        dr.font.color.rgb = RGBColor(75, 85, 99)

        col_top[col_idx] = current_top + Inches(1.0)


def describe_component(props):
    if isinstance(props, str):
        return props
    if not isinstance(props, dict):
        return ""
    if "description" in props:
        return str(props["description"])
    parts = []
    for k in ("backgroundColor", "textColor", "rounded", "padding", "border", "typography"):
        if k in props:
            parts.append(f"{k}: {props[k]}")
    return " · ".join(parts)


def add_dos_donts_slide(prs, body):
    dos, donts = extract_dos_donts(body)
    if not dos and not donts:
        return

    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(6.0), Inches(0.6),
                 "Do", size=28, bold=True, color=(16, 122, 60))
    add_text_box(s, Inches(6.9), Inches(0.5), Inches(6.0), Inches(0.6),
                 "Don't", size=28, bold=True, color=(180, 35, 24))

    def add_bullets(left, items):
        if not items:
            return
        box = s.shapes.add_textbox(left, Inches(1.4), Inches(6.0), Inches(5.7))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in items[:7]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = "• " + strip_inline_md(item)
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(8)

    add_bullets(Inches(0.6), dos)
    add_bullets(Inches(6.9), donts)


def add_reconstruction_slide(prs, body):
    section5 = extract_section(body, 5)
    if not section5:
        return

    s = add_blank_slide(prs)
    add_text_box(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
                 "Reconstruction notes", size=28, bold=True, color=(23, 23, 23))

    stack = extract_subsection(section5, "Suggested stack")
    quick_wins = extract_subsection_bullets(section5, "Quick wins")
    tricky = extract_subsection_bullets(section5, "Tricky bits")

    top = Inches(1.5)
    if stack:
        add_text_box(s, Inches(0.6), top, Inches(2.0), Inches(0.4),
                     "Stack", size=14, bold=True, color=(23, 23, 23))
        add_text_box(s, Inches(2.6), top, Inches(10.0), Inches(0.6),
                     stack, size=12, color=(75, 85, 99))
        top += Inches(1.0)

    if quick_wins:
        add_text_box(s, Inches(0.6), top, Inches(5.5), Inches(0.4),
                     "Quick wins", size=14, bold=True, color=(16, 122, 60))
        wb = s.shapes.add_textbox(Inches(0.6), top + Inches(0.4), Inches(5.5), Inches(4.0))
        wf = wb.text_frame
        wf.word_wrap = True
        first = True
        for item in quick_wins[:6]:
            p = wf.paragraphs[0] if first else wf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = "• " + strip_inline_md(item)
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(75, 85, 99)
            p.space_after = Pt(6)

    if tricky:
        add_text_box(s, Inches(6.9), top, Inches(5.5), Inches(0.4),
                     "Tricky bits", size=14, bold=True, color=(180, 110, 24))
        tb = s.shapes.add_textbox(Inches(6.9), top + Inches(0.4), Inches(5.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for item in tricky[:6]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = "• " + strip_inline_md(item)
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(75, 85, 99)
            p.space_after = Pt(6)


def extract_subsection(section_text, heading):
    pattern = rf"###\s+{re.escape(heading)}\s*\n(.*?)(?=###\s+|\Z)"
    m = re.search(pattern, section_text, re.DOTALL)
    if not m:
        return ""
    text = m.group(1).strip()
    text = re.sub(r"\n+", " ", text).strip()
    return strip_inline_md(text)[:300]


def extract_subsection_bullets(section_text, heading):
    pattern = rf"###\s+{re.escape(heading)}\s*\n(.*?)(?=###\s+|\Z)"
    m = re.search(pattern, section_text, re.DOTALL)
    if not m:
        return []
    bullets = []
    for line in m.group(1).split("\n"):
        stripped = line.strip()
        if re.match(r"^[-*]\s+", stripped):
            bullets.append(re.sub(r"^[-*]\s+", "", stripped))
    return bullets


# =============================================================================
# README-claude-design.md
# =============================================================================

def emit_readme(out_dir, source_name, emitted):
    has_pptx = "brand-kit.pptx" in emitted
    has_docx = "brand-overview.docx" in emitted
    has_css = "tokens.css" in emitted
    has_tw = "tailwind.config.ts" in emitted

    lines = [
        f"# Claude Design upload bundle — {source_name}",
        "",
        "This bundle was generated by anydesign for use with **Claude Design**",
        "(claude.ai/design). Each file targets one of Claude Design's setup inputs.",
        "",
        "## Files in this bundle",
        "",
    ]
    if has_pptx:
        lines += [
            "- **`brand-kit.pptx`** — primary asset. A multi-slide deck covering atmosphere, palette, typography, components, and Do's/Don'ts. Claude Design ingests PPTX directly.",
        ]
    if has_docx:
        lines += [
            "- **`brand-overview.docx`** — the full `design.md` rendered as a Word document. Use it as the brand brief.",
        ]
    if has_css or has_tw:
        repo_files = []
        if has_css:
            repo_files.append("`tokens.css`")
        if has_tw:
            repo_files.append("`tailwind.config.ts`")
        lines += [
            f"- **{', '.join(repo_files)}** — drop into an empty repo (or your design-tokens repo) and link the repo to Claude Design. Its system extractor reads CSS custom properties and Tailwind configs.",
        ]
    lines += [
        "",
        "## How to upload",
        "",
        "1. Open **claude.ai/design** and create a new project (or open Settings → Design system).",
        "2. In the **brand and product assets** section, upload:",
    ]
    if has_pptx:
        lines.append("   - `brand-kit.pptx`")
    if has_docx:
        lines.append("   - `brand-overview.docx`")
    if has_css or has_tw:
        lines.append("   - For the code-repo path: push `tokens.css` and/or `tailwind.config.ts` to a small repo and link it under **Code repository**.")
    lines += [
        "",
        "3. Let Claude Design extract the system. Review the detected colors, typography, and components — adjust where needed.",
        "4. Future projects in this organization will default to this system.",
        "",
        "## Troubleshooting",
        "",
        "- If Claude Design's extracted palette differs from your tokens, re-upload the PPTX first (richest visual context) and use the **Remix** button.",
        "- If colors come through but typography doesn't, the font family in the bundle is a fallback chain — name the exact font in the design-system setup notes.",
        "- **Chinese (CJK) text renders in the wrong font / shows non-Simplified glyphs:** the cloud render environment may not ship the CJK fonts in the fallback chain. The bundle's font stacks already prioritize `Noto Sans SC` (cloud-friendly) ahead of local-only fonts like `PingFang SC`. To guarantee it, load Noto Sans SC explicitly in the generated HTML `<head>`:",
        "",
        "  ```html",
        '  <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@400;500;700&display=swap" rel="stylesheet">',
        "  ```",
        "",
        "- The DTCG `design-tokens.json` from anydesign is the canonical source. If Claude Design ever supports DTCG natively, you can drop it in directly.",
        "",
        "---",
        "",
        "*Generated by [anydesign](https://github.com/uxKero/anydesign) · `scripts/export_for_claude_design.py`*",
        "",
    ]
    (out_dir / "README-claude-design.md").write_text("\n".join(lines), encoding="utf-8")
    return True


# =============================================================================
# Main
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Export a design.md + design-tokens.json bundle for Claude Design.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("design_md", help="Path to design.md")
    parser.add_argument("tokens_json", nargs="?", default=None,
                        help="Path to design-tokens.json (optional; auto-detected as sibling file if omitted)")
    parser.add_argument("--out", default="claude-design-bundle",
                        help="Output directory (default: claude-design-bundle)")
    parser.add_argument("--skip", default="",
                        help=f"Comma-separated formats to skip: {','.join(ALL_FORMATS)}")
    parser.add_argument("--no-cjk-fallback", action="store_true",
                        help="Don't append the Simplified-Chinese font fallback chain to exported font stacks")
    args = parser.parse_args()

    design_path = Path(args.design_md)
    if not design_path.exists():
        print(f"File not found: {design_path}", file=sys.stderr)
        sys.exit(2)

    if args.tokens_json:
        tokens_path = Path(args.tokens_json)
    else:
        sibling = design_path.parent / "design-tokens.json"
        tokens_path = sibling if sibling.exists() else None

    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    skipped = {f.strip() for f in args.skip.split(",") if f.strip()}
    unknown = skipped - set(ALL_FORMATS)
    if unknown:
        print(f"warn: unknown --skip values ignored: {', '.join(sorted(unknown))}", file=sys.stderr)
    formats = [f for f in ALL_FORMATS if f not in skipped]

    text = design_path.read_text(encoding="utf-8")
    fm_text, body = split_frontmatter(text)

    fm = None
    if fm_text:
        if HAS_YAML:
            fm = parse_frontmatter(fm_text)
        else:
            print("warn: pyyaml not installed — frontmatter parsing skipped. Install with: pip install pyyaml", file=sys.stderr)

    tokens = None
    if tokens_path and tokens_path.exists():
        try:
            tokens = json.loads(tokens_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as e:
            print(f"warn: failed to parse {tokens_path}: {e}", file=sys.stderr)

    source_name = (fm or {}).get("name") or design_path.stem

    print(f"anydesign → Claude Design bundle")
    print(f"  design.md:    {design_path}")
    print(f"  tokens.json:  {tokens_path if tokens_path else '(none — CSS/Tailwind disabled)'}")
    print(f"  output dir:   {out_dir}")
    print()

    emitted = []

    if "css" in formats:
        if tokens:
            ok = emit_tokens_css(tokens, out_dir / "tokens.css", source_name,
                                 cjk_fallback=not args.no_cjk_fallback)
            print(f"  {'✓' if ok else '✗'} tokens.css")
            if ok:
                emitted.append("tokens.css")
        else:
            print("  - tokens.css (skipped: no design-tokens.json)")

    if "tailwind" in formats:
        if tokens:
            ok = emit_tailwind_config(tokens, out_dir / "tailwind.config.ts", source_name,
                                      cjk_fallback=not args.no_cjk_fallback)
            print(f"  {'✓' if ok else '✗'} tailwind.config.ts")
            if ok:
                emitted.append("tailwind.config.ts")
        else:
            print("  - tailwind.config.ts (skipped: no design-tokens.json)")

    if "docx" in formats:
        if HAS_DOCX:
            ok = emit_brand_overview_docx(fm, body, out_dir / "brand-overview.docx", source_name)
            print(f"  {'✓' if ok else '✗'} brand-overview.docx")
            if ok:
                emitted.append("brand-overview.docx")
        else:
            print("  - brand-overview.docx (skipped: install python-docx)")

    if "pptx" in formats:
        if HAS_PPTX:
            ok = emit_brand_kit_pptx(fm, tokens, body, out_dir / "brand-kit.pptx", source_name)
            print(f"  {'✓' if ok else '✗'} brand-kit.pptx")
            if ok:
                emitted.append("brand-kit.pptx")
        else:
            print("  - brand-kit.pptx (skipped: install python-pptx)")

    if "readme" in formats:
        emit_readme(out_dir, source_name, emitted)
        print(f"  ✓ README-claude-design.md")
        emitted.append("README-claude-design.md")

    print()
    print(f"Done. {len(emitted)} file(s) in {out_dir.resolve()}")
    if not emitted:
        sys.exit(1)


if __name__ == "__main__":
    main()
